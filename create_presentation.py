#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Erstellt eine PowerPoint-Präsentation für die Data Mining Aufgabe zur Betrugserkennung
Strukturiert nach CRISP-DM-Prozess
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Erstellt eine Titelfolie"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def create_content_slide(prs, title, bullet_points=None):
    """Erstellt eine Folie mit Aufzählungspunkten"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    if bullet_points:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for i, (level, text) in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.level = level
    
    return slide

def create_table_slide(prs, title, data, col_widths=None):
    """Erstellt eine Folie mit Tabelle"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    
    rows = len(data)
    cols = len(data[0])
    
    left = Inches(1.0)
    top = Inches(2.0)
    width = Inches(8.0)
    height = Inches(0.8 * rows)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Spaltenbreiten setzen
    if col_widths:
        for i, width in enumerate(col_widths):
            table.columns[i].width = Inches(width)
    
    # Tabelle füllen
    for i, row in enumerate(data):
        for j, cell_text in enumerate(row):
            cell = table.rows[i].cells[j]
            cell.text = str(cell_text)
            
            # Header-Zeile formatieren
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        run.font.size = Pt(12)
            else:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(11)
    
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Folie 1: Titel
    create_title_slide(
        prs,
        "Data Mining: Betrugserkennung",
        "Klassifikation von E-Commerce Bestellungen"
    )
    
    # Folie 2: Agenda
    create_content_slide(
        prs,
        "Agenda",
        [
            (0, "1. Einführung und Problemstellung"),
            (0, "2. Datenanalyse und -vorbereitung (CRISP-DM)"),
            (0, "3. Attributanalyse und Feature Engineering"),
            (0, "4. Modellauswahl und Training"),
            (0, "5. Modellbewertung und Optimierung"),
            (0, "6. Vorhersage und Fazit")
        ]
    )
    
    # Folie 3: Problemstellung
    create_content_slide(
        prs,
        "1. Einführung und Problemstellung",
        [
            (0, "Szenario: Betrugserkennung im Online-Handel"),
            (1, "Herausforderung: Ware gegen Geld nicht direkt umsetzbar"),
            (1, "Risiko: Bestellungen ohne Zahlungseingang"),
            (0, "Data Mining Aufgabe:"),
            (1, "Typ: Klassifikationsproblem (Überwachtes Lernen)"),
            (1, "Zielattribut: TARGET_BETRUG (diskret: ja/nein)"),
            (0, "Datensatz: 30.000 E-Commerce Bestellungen"),
            (1, "Training: 24.000 Datensätze (80% - Holdout-Methode)"),
            (1, "Test: 6.000 Datensätze (20%)"),
            (1, "Klassifizierung: 20.000 neue Bestellungen")
        ]
    )
    
    # Folie 4: CRISP-DM und Data Understanding
    create_content_slide(
        prs,
        "2. Datenanalyse (CRISP-DM: Data Understanding)",
        [
            (0, "Explorative Datenanalyse der Trainingsdaten:"),
            (1, "30.000 Bestellungen mit 43 Attributen"),
            (1, "Zielattribut TARGET_BETRUG analysiert"),
            (0, "Zentrale Erkenntnis: Stark unbalancierte Klassen"),
            (1, "Nur 1.746 Betrugsfälle (5,82%)"),
            (1, "28.254 legitime Bestellungen (94,18%)"),
            (0, "Implikation für Data Mining:"),
            (1, "❌ Accuracy allein als Gütemaß ungeeignet"),
            (1, "✓ Precision & Recall sind entscheidend"),
            (1, "Modell könnte \"immer nein\" vorhersagen (94% Accuracy!)"),
            (1, "→ Betrugsfälle würden nicht erkannt werden")
        ]
    )
    
    # Folie 5: Feature Engineering
    create_content_slide(
        prs,
        "3. Attributanalyse & Feature Engineering (Data Preparation)",
        [
            (0, "Analyse der Artikelnummer-Attribute (ANUMMER_01 bis ANUMMER_10):"),
            (1, "Repräsentieren bestellte Artikel (Artikelnummern)"),
            (1, "Datentyp: Diskret (kategorisch)"),
            (0, "Identifizierte Probleme:"),
            (1, "❌ Limitation: Maximal 10 Artikel erfassbar"),
            (1, "❌ Bei >10 Artikeln: Datenverlust"),
            (1, "❌ ANUMMER_02 bis ANUMMER_10 größtenteils NULL"),
            (1, "❌ Hohe Dimensionalität mit vielen Leereinträgen"),
            (0, "Feature Engineering - Optimierungsansatz:"),
            (1, "✓ Zusammenführung zu ANUMMER_LIST (aggregiert)"),
            (1, "✓ Reduzierung der Dimensionalität (10→1 Attribut)"),
            (1, "✓ Bessere Datenqualität für ML-Algorithmen")
        ]
    )
    
    # Folie 6: Modellierung
    create_content_slide(
        prs,
        "4. Modellauswahl und Training (CRISP-DM: Modeling)",
        [
            (0, "Implementierung: KNIME Analytics Platform"),
            (0, "Trainingsmethode: Holdout-Methode"),
            (1, "80% Trainingsdaten (24.000 Datensätze)"),
            (1, "20% Testdaten (6.000 Datensätze)"),
            (0, "Drei Klassifikationsverfahren (Überwachtes Lernen):"),
            (1, "1. Decision Tree (Entscheidungsbaum, Gini Index)"),
            (2, "Trennscharfe Entscheidungsregeln"),
            (1, "2. Logistic Regression (Logistische Regression)"),
            (2, "Lineare Trennung der Klassen"),
            (1, "3. Naive Bayes Classifier"),
            (2, "Probabilistischer Ansatz, bedingte Wahrscheinlichkeiten")
        ]
    )
    
    # Folie 7: Konfusionsmatrix Erklärung
    create_content_slide(
        prs,
        "5. Modellbewertung (CRISP-DM: Evaluation)",
        [
            (0, "Bewertung mit Konfusionsmatrix:"),
            (1, "True Positive (TP): Betrug korrekt als Betrug erkannt"),
            (1, "True Negative (TN): Kein Betrug korrekt als legitim erkannt"),
            (1, "False Positive (FP): Legitim fälschlich als Betrug erkannt"),
            (1, "False Negative (FN): Betrug fälschlich als legitim erkannt"),
            (0, "Gütekriterien:"),
            (1, "Accuracy = (TP+TN) / Gesamt"),
            (2, "Gesamtgenauigkeit aller Vorhersagen"),
            (1, "Precision = TP / (TP+FP)"),
            (2, "Wie viele Betrugsvorhersagen waren korrekt?"),
            (1, "Recall = TP / (TP+FN)"),
            (2, "Wie viele Betrugsfälle wurden erkannt?")
        ]
    )
    
    # Folie 8: Decision Tree Confusion Matrix
    table_data = [
        ["", "Predicted: nein", "Predicted: ja"],
        ["Actual: nein", "5466 (TN)", "193 (FP)"],
        ["Actual: ja", "300 (FN)", "36 (TP)"]
    ]
    
    create_table_slide(
        prs,
        "5a. Decision Tree - Confusion Matrix",
        table_data,
        col_widths=[3.0, 2.5, 2.5]
    )
    
    # Folie 9: Decision Tree Metriken
    create_content_slide(
        prs,
        "5a. Decision Tree - Metriken",
        [
            (0, "Berechnete Gütekriterien:"),
            (1, "Accuracy: 91,78% = (5466+36) / 6000"),
            (1, "Precision: 15,72% = 36 / (36+193)"),
            (1, "Recall: 9,33% = 36 / (36+300)"),
            (0, "Interpretation:"),
            (1, "✓ Erkennt tatsächlich Betrugsfälle (36 TP)"),
            (1, "✓ Besser als \"immer nein\" Baseline"),
            (1, "⚠ Hohe False-Positive Rate (193 Fehlalarme)"),
            (1, "⚠ Niedriger Recall (300 FN, 89% nicht erkannt)"),
            (0, "Bewertung: Einziges Modell mit relevantem Recall")
        ]
    )
    
    # Folie 10: Logistic Regression
    create_content_slide(
        prs,
        "5b. Logistic Regression - Ergebnisse",
        [
            (0, "Confusion Matrix:"),
            (1, "TN: 5666, FP: 0, FN: 336, TP: 0"),
            (1, "Alle Vorhersagen: \"nein\" (kein Betrug)"),
            (0, "Berechnete Metriken:"),
            (1, "Accuracy: 94,40% = 5666 / 6000"),
            (1, "Precision: 0% (keine Betrugsvorhersagen)"),
            (1, "Recall: 0% = 0 / 336 (alle Betrugsfälle übersehen)"),
            (0, "❌ Fazit: Modell unbrauchbar für Betrugserkennung"),
            (1, "Paradox: Hohe Accuracy durch Class Imbalance"),
            (1, "Modell hat aus Unbalance gelernt, immer \"nein\" zu sagen"),
            (1, "Kein einziger Betrugsfall erkannt (0 TP)")
        ]
    )
    
    # Folie 11: Naive Bayes
    create_content_slide(
        prs,
        "5c. Naive Bayes - Ergebnisse",
        [
            (0, "Confusion Matrix:"),
            (1, "TN: 5663, FP: 3, FN: 333, TP: 3"),
            (0, "Berechnete Metriken:"),
            (1, "Accuracy: 94,44% = (5663+3) / 6000"),
            (1, "Precision: 50,00% = 3 / (3+3)"),
            (1, "Recall: 0,89% = 3 / (3+333)"),
            (0, "Bewertung:"),
            (1, "✓ Hohe Precision bei Betrugsvorhersagen (50%)"),
            (1, "✓ Sehr wenig Fehlalarme (nur 3 FP)"),
            (1, "❌ Extrem niedriger Recall (0,89%)"),
            (1, "❌ Erkennt fast keine Betrugsfälle (nur 3 von 336)"),
            (0, "Fazit: Zu konservativ, praktisch unbrauchbar")
        ]
    )
    
    # Folie 12: Modellvergleich
    comparison_data = [
        ["Modell", "Accuracy", "Precision", "Recall", "TP", "Bewertung"],
        ["Decision Tree", "91,78%", "15,72%", "9,33%", "36", "⭐ GEWÄHLT"],
        ["Logistic Regression", "94,40%", "0%", "0%", "0", "❌ Unbrauchbar"],
        ["Naive Bayes", "94,44%", "50,00%", "0,89%", "3", "⚠ Zu konservativ"]
    ]
    
    create_table_slide(
        prs,
        "5d. Vergleichende Modellbewertung",
        comparison_data,
        col_widths=[2.3, 1.3, 1.3, 1.3, 0.8, 2.0]
    )
    
    # Folie 13: Optimierung
    create_content_slide(
        prs,
        "5e. Maßnahmen zur Ergebnisverbesserung",
        [
            (0, "Hauptproblem: Class Imbalance (5,82% Betrug)"),
            (0, "Resampling-Techniken:"),
            (1, "Undersampling: Reduzierung der Mehrheitsklasse"),
            (2, "Zufällige Auswahl von legitimen Bestellungen"),
            (1, "SMOTE (Synthetic Minority Over-sampling Technique)"),
            (2, "Synthetische Generierung zusätzlicher Betrugsfälle"),
            (0, "Algorithmus-Optimierung:"),
            (1, "Hyperparameter-Tuning (Complexity_Penalty, Minimum_Support)"),
            (1, "Cost-Sensitive Learning (höhere Kosten für FN)"),
            (0, "Weitere Ansätze:"),
            (1, "Ensemble-Methoden (Random Forest, XGBoost)"),
            (1, "Feature Engineering (ANUMMER_LIST verwenden)")
        ]
    )
    
    # Folie 14: Vorhersage und Deployment
    create_content_slide(
        prs,
        "6. Vorhersage (CRISP-DM: Deployment)",
        [
            (0, "Gewähltes Modell für Produktivbetrieb:"),
            (1, "✓ Decision Tree (Entscheidungsbaum)"),
            (0, "Begründung der Auswahl:"),
            (1, "Einziges Modell mit relevantem Recall (9,33%)"),
            (1, "Erkennt tatsächlich Betrugsfälle (36 TP)"),
            (1, "Trade-off: Niedrigere Accuracy, aber Betrugserkennnung"),
            (0, "Anwendung auf Klassifizierungsdaten:"),
            (1, "20.000 neue Bestellungen klassifiziert"),
            (1, "Datei: Klassifizierungsdaten-tree-predictions.csv"),
            (0, "Kritische Einschätzung:"),
            (1, "Performance noch verbesserungswürdig"),
            (1, "Optimierungsmaßnahmen empfohlen (Resampling, etc.)")
        ]
    )
    
    # Folie 15: Zusammenfassung
    create_content_slide(
        prs,
        "Zusammenfassung und Ausblick",
        [
            (0, "Zentrale Erkenntnisse:"),
            (1, "Class Imbalance größte Herausforderung bei Betrugserkennung"),
            (1, "Accuracy irreführend bei unbalancierten Daten"),
            (1, "Precision & Recall entscheidende Gütekriterien"),
            (1, "CRISP-DM-Prozess systematisch durchgeführt"),
            (0, "Durchgeführte Schritte:"),
            (1, "✓ Data Understanding (Class Imbalance identifiziert)"),
            (1, "✓ Data Preparation (ANUMMER Feature Engineering)"),
            (1, "✓ Modeling (3 Klassifikationsverfahren trainiert)"),
            (1, "✓ Evaluation (Konfusionsmatrix, Metriken)"),
            (1, "✓ Deployment (Decision Tree für Vorhersage gewählt)"),
            (0, "Vielen Dank für Ihre Aufmerksamkeit!")
        ]
    )
    
    # Präsentation speichern
    output_file = "Data_Mining_Betrugserkennung_Praesentation.pptx"
    prs.save(output_file)
    print(f"✓ Präsentation erfolgreich erstellt: {output_file}")
    print(f"✓ Anzahl Folien: {len(prs.slides)}")
    print(f"\nStruktur nach CRISP-DM:")
    print("  1. Einführung und Problemstellung")
    print("  2. Data Understanding (Class Imbalance)")
    print("  3. Data Preparation (Feature Engineering)")
    print("  4. Modeling (3 Klassifikationsverfahren)")
    print("  5. Evaluation (Konfusionsmatrix, Metriken)")
    print("  6. Deployment (Vorhersage)")

if __name__ == "__main__":
    main()
